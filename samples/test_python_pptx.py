from pptx import Presentation
from pptx.util import Inches, Pt

prs = Presentation()

# Title slide
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Markdown 幻灯片工具对比"
subtitle.text = "Generated with python-pptx"

# Content slide
bullet_slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(bullet_slide_layout)
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = "Why Markdown for Slides?"
tf = body_shape.text_frame
tf.text = "Version Control: Plain text works with Git"
p = tf.add_paragraph()
p.text = "Portability: Write once, export anywhere"
p.level = 0
p = tf.add_paragraph()
p.text = "Focus on Content: Separate content from styling"
p.level = 0

prs.save('/home/zliang/workspace/omscs/HongLab/slides_marker/output/python_pptx_output.pptx')
print("PPTX created successfully!")
